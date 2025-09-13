import streamlit as st
import os
import tempfile
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
import base64
from io import BytesIO
import time
from PIL import Image
import pandas as pd
import zipfile
import subprocess
import sys
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from pptx.util import Inches

# Set page configuration
st.set_page_config(
    page_title="Advanced Document Converter",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for styling - Fixed syntax errors
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
        padding: 1rem;
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
        border-radius: 10px;
    }
    .sub-header {
        font-size: 1.8rem;
        color: #2c3e;
        margin-bottom: 1rem;
        padding: 0.5rem;
        border-left: 5px solid #1f77b4;
    }
    .converter-option {
        background-color: #f8f9fa;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        border: 1px solid #e9ecef;
    }
    .upload-box {
        border: 2px dashed #6c757d;
        padding: 2rem;
        border-radius: 10px;
        text-align: center;
        margin-bottom: 2rem;
        background-color: #f9f9f9;
    }
    .success-box {
        background-color: #d4edda;
        color: #155724;
        padding: 15px;
        border-radius: 5px;
        margin: 10px 0;
    }
    .info-box {
        background-color: #d1ecf1;
        color: #0c5460;
        padding: 15px;
        border-radius: 5px;
        margin: 10px 0;
    }
    .stButton>button {
        width: 100%;
        border-radius: 5px;
        height: 3rem;
        font-weight: bold;
        background: linear-gradient(135deg, #6a11cb 0%, #2575fc 100%);
        color: white;
    }
    .footer {
        text-align: center;
        margin-top: 2rem;
        padding: 1rem;
        color: #6c757d;
        background-color: #f8f9fa;
        border-radius: 10px;
    }
    .sidebar .sidebar-content {
        background: linear-gradient(180deg, #4b6cb7 0%, #182848 100%);
        color: white;
    }
    .sidebar .sidebar-content .stRadio label {
        color: white;
    }
</style>
""", unsafe_allow_html=True)

# Create necessary directories
os.makedirs("temp", exist_ok=True)
os.makedirs("output", exist_ok=True)

def save_uploaded_file(uploaded_file, directory="temp"):
    """Save uploaded file to directory and return file path"""
    try:
        file_path = os.path.join(directory, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        return file_path
    except Exception as e:
        st.error(f"Error saving file: {str(e)}")
        return None

def pdf_to_word(pdf_file):
    """Convert PDF file to Word document"""
    try:
        # Create a new Word document
        doc = Document()
        
        # Open the PDF file
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Process each page
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Extract text
            text = page.get_text()
            
            # Add text to Word document
            if text.strip():
                doc.add_paragraph(text)
        
        pdf_document.close()
        
        # Save to BytesIO buffer
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        return buffer
        
    except Exception as e:
        st.error(f"Error converting PDF: {str(e)}")
        return None

def word_to_pdf(docx_file):
    """Convert Word document to PDF using LibreOffice (cross-platform)"""
    try:
        # Save the uploaded file
        docx_path = save_uploaded_file(docx_file)
        if not docx_path:
            return None
        
        # Create output path
        pdf_path = docx_path.replace('.docx', '.pdf').replace('.doc', '.pdf')
        
        # Convert using LibreOffice
        if sys.platform == "win32":
            # Windows
            cmd = f'soffice --headless --convert-to pdf "{docx_path}" --outdir "{os.path.dirname(docx_path)}"'
        else:
            # Linux/Mac
            cmd = f'libreoffice --headless --convert-to pdf "{docx_path}" --outdir "{os.path.dirname(docx_path)}"'
        
        subprocess.run(cmd, shell=True, check=True)
        
        # Read the generated PDF
        with open(pdf_path, 'rb') as f:
            pdf_data = f.read()
        
        # Clean up temporary files
        os.unlink(docx_path)
        if os.path.exists(pdf_path):
            os.unlink(pdf_path)
        
        return BytesIO(pdf_data)
        
    except Exception as e:
        st.error(f"Error converting Word to PDF: {str(e)}")
        st.info("Please make sure LibreOffice is installed on your system.")
        return None

def pdf_to_pptx(pdf_file):
    """Convert PDF to PowerPoint presentation"""
    try:
        # Open the PDF file
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PDF Conversion"
        subtitle.text = f"Converted from {pdf_file.name}"
        
        # Process each page and add as a new slide
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Extract text
            text = page.get_text()
            
            # Add a new slide for each page
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add text box
            left = Inches(0.5)
            top = Inches(1)
            width = Inches(9)
            height = Inches(6)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = text
        
        pdf_document.close()
        
        # Save to BytesIO buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return buffer
        
    except Exception as e:
        st.error(f"Error converting PDF to PowerPoint: {str(e)}")
        return None

def pdf_to_jpg(pdf_file, page_number=0):
    """Convert PDF page to JPG image"""
    try:
        # Open the PDF file
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Get the specified page
        if page_number >= len(pdf_document):
            page_number = 0
            
        page = pdf_document.load_page(page_number)
        
        # Convert to image
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("jpeg")
        
        pdf_document.close()
        
        return BytesIO(img_data)
        
    except Exception as e:
        st.error(f"Error converting PDF to JPG: {str(e)}")
        return None

def jpg_to_pdf(image_files):
    """Convert JPG images to PDF"""
    try:
        # Create a new PDF
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        
        for img_file in image_files:
            # Open image
            img = Image.open(img_file)
            img.thumbnail((width - 100, height - 100))  # Resize to fit page
            
            # Convert to PDF page
            img_buffer = BytesIO()
            img.save(img_buffer, format='JPEG')
            img_buffer.seek(0)
            
            # Add to PDF
            c.drawImage(ImageReader(img_buffer), 50, 50, width=img.width, height=img.height)
            c.showPage()
        
        c.save()
        buffer.seek(0)
        
        return buffer
        
    except Exception as e:
        st.error(f"Error converting JPG to PDF: {str(e)}")
        return None

def pdf_to_excel(pdf_file):
    """Convert PDF to Excel"""
    try:
        # Open the PDF file
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Extract text
        text_content = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text_content += page.get_text() + "\n"
        
        pdf_document.close()
        
        # Create a DataFrame with the text
        df = pd.DataFrame({"Extracted Text": text_content.split("\n")})
        
        # Save to BytesIO buffer
        buffer = BytesIO()
        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
            df.to_excel(writer, index=False)
        buffer.seek(0)
        
        return buffer
        
    except Exception as e:
        st.error(f"Error converting PDF to Excel: {str(e)}")
        return None





# def excel_to_pdf(excel_file):
#     """Convert Excel to PDF using LibreOffice (no hardcoded path)"""
#     try:
#         # Save uploaded Excel file
#         excel_path = save_uploaded_file(excel_file)
#         if not excel_path:
#             return None

#         # Output PDF path
#         pdf_path = excel_path.rsplit(".", 1)[0] + ".pdf"

#         # Choose command based on OS
#         cmd = (
#             f'soffice --headless --convert-to pdf:calc_pdf_Export "{excel_path}" --outdir "{os.path.dirname(excel_path)}"'
#             if sys.platform == "win32"
#             else f'libreoffice --headless --convert-to pdf:calc_pdf_Export "{excel_path}" --outdir "{os.path.dirname(excel_path)}"'
#         )

#         # Run command and capture logs
#         result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
#         if result.returncode != 0:
#             raise Exception(f"LibreOffice error:\n{result.stderr}")

#         # Read PDF
#         with open(pdf_path, "rb") as f:
#             pdf_data = f.read()

#         # Cleanup
#         os.unlink(excel_path)
#         os.unlink(pdf_path)

#         return BytesIO(pdf_data)

#     except Exception as e:
#         st.error(f"Error converting Excel to PDF: {str(e)}")
#         st.info("Make sure LibreOffice is installed and added to your PATH.")
#         return None









def merge_pdfs(pdf_files):
    """Merge multiple PDF files into one"""
    try:
        # Create a new PDF document
        merged_pdf = fitz.open()
        
        # Append each PDF
        for pdf_file in pdf_files:
            pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
            merged_pdf.insert_pdf(pdf_document)
            pdf_document.close()
        
        # Save to bytes buffer
        buffer = BytesIO()
        merged_pdf.save(buffer)
        buffer.seek(0)
        merged_pdf.close()
        
        return buffer
        
    except Exception as e:
        st.error(f"Error merging PDFs: {str(e)}")
        return None

def split_pdf(pdf_file, split_type="single", start_page=1, end_page=1):
    """Split a PDF file into multiple files"""
    try:
        # Open the PDF file
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        if split_type == "single":
            # Create a zip file containing all pages as separate PDFs
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for page_num in range(len(pdf_document)):
                    # Create a new PDF for each page
                    single_page_pdf = fitz.open()
                    single_page_pdf.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
                    
                    # Save to bytes buffer
                    page_buffer = BytesIO()
                    single_page_pdf.save(page_buffer)
                    page_buffer.seek(0)
                    
                    # Add to zip
                    zip_file.writestr(f"page_{page_num + 1}.pdf", page_buffer.getvalue())
                    single_page_pdf.close()
            
            pdf_document.close()
            zip_buffer.seek(0)
            return zip_buffer
            
        else:  # Page range
            # Create a new PDF with the specified page range
            extracted_pdf = fitz.open()
            extracted_pdf.insert_pdf(pdf_document, from_page=start_page-1, to_page=end_page-1)
            
            # Save to bytes buffer
            buffer = BytesIO()
            extracted_pdf.save(buffer)
            buffer.seek(0)
            
            pdf_document.close()
            extracted_pdf.close()
            
            return buffer
            
    except Exception as e:
        st.error(f"Error splitting PDF: {str(e)}")
        return None
    





def compress_pdf(pdf_file, compression_level=3):
    """Compress a PDF file by recompressing images and optimizing storage"""
    try:
        # Open the PDF
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")

        # Map compression level to image quality
        quality_map = {
            1: 95,  # light
            2: 85,  # medium
            3: 75,  # balanced
            4: 65,  # strong
            5: 50   # max
        }
        image_quality = quality_map.get(compression_level, 75)

        # Process each page and its images
        for page_index in range(len(pdf_document)):
            page = pdf_document[page_index]
            images = page.get_images(full=True)

            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]

                # Open with Pillow
                image = Image.open(BytesIO(image_bytes)).convert("RGB")

                # Recompress with selected quality
                img_buffer = BytesIO()
                image.save(img_buffer, format="JPEG", quality=image_quality, optimize=True)
                new_image_bytes = img_buffer.getvalue()

                # Replace old image with new compressed one
                pdf_document.update_stream(xref, new_image_bytes)

        # Save with cleanup & compression
        buffer = BytesIO()
        pdf_document.save(
            buffer,
            deflate=True,      # compress streams
            garbage=4,         # remove unused objects
            clean=True,        # clean up
            incremental=False
        )
        buffer.seek(0)
        pdf_document.close()

        return buffer

    except Exception as e:
        st.error(f"Error compressing PDF: {str(e)}")
        return None  
st.sidebar.header("Made By Subhadip 😎")
def main():
    # Header
    st.markdown('<h1 class="main-header">📄 Advanced Document Converter</h1>', unsafe_allow_html=True)
    
    # Sidebar
    with st.sidebar:
        st.header("🔧 Conversion Tools")
        
        # Tool selection
        tool_option = st.radio(
            "Select Conversion Type:",
            [
                "PDF to Word", 
                "Word to PDF", 
                "Merge PDFs", 
                "Split PDF", 
                "Compress PDF",
                "PDF to PowerPoint",
                "PDF to JPG", 
                "JPG to PDF",
                "PDF to Excel",
                # "Excel to PDF"
            ]
        )
        
        st.markdown("---")
        st.header("ℹ️ About")
        st.info("""
        This tool allows you to convert between various document formats.
        
        **Note:** Some conversions require LibreOffice to be installed on your system.
        """)
        
        st.markdown("---")
        st.header("📝 Instructions")
        st.write("""
        1. Select conversion type
        2. Upload your file(s)
        3. Configure options (if any)
        4. Click the convert button
        5. Download the converted file
        """)
    
    # Main content based on selected tool
    if tool_option == "PDF to Word":
        st.markdown('<h2 class="sub-header">PDF to Word Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=["pdf"],
            help="Select a PDF file to convert to Word format"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Convert button
            if st.button("Convert PDF to Word"):
                with st.spinner("Converting PDF to Word..."):
                    # Convert PDF to Word
                    word_data = pdf_to_word(uploaded_file)
                    
                    if word_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = Path(uploaded_file.name).stem + ".docx"
                        st.download_button(
                            label="📥 Download Word Document",
                            data=word_data,
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "Word to PDF":
        st.markdown('<h2 class="sub-header">Word to PDF Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a Word document",
            type=["docx", "doc"],
            help="Select a Word file to convert to PDF format"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Convert button
            if st.button("Convert Word to PDF"):
                with st.spinner("Converting Word to PDF..."):
                    # Convert Word to PDF
                    pdf_data = word_to_pdf(uploaded_file)
                    
                    if pdf_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = Path(uploaded_file.name).stem + ".pdf"
                        st.download_button(
                            label="📥 Download PDF Document",
                            data=pdf_data,
                            file_name=output_filename,
                            mime="application/pdf"
                        )
                    else:
                        st.info("💡 If conversion fails, please ensure LibreOffice is installed on your system.")
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "Merge PDFs":
        st.markdown('<h2 class="sub-header">Merge PDF Files</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_files = st.file_uploader(
            "Choose PDF files to merge",
            type=["pdf"],
            help="Select multiple PDF files to merge",
            accept_multiple_files=True
        )
        
        if uploaded_files and len(uploaded_files) > 1:
            st.write(f"**Files to merge:** {len(uploaded_files)} files selected")
            
            # Convert button
            if st.button("Merge PDFs"):
                with st.spinner("Merging PDF files..."):
                    # Merge PDFs
                    merged_pdf = merge_pdfs(uploaded_files)
                    
                    if merged_pdf:
                        st.markdown('<div class="success-box">✅ PDFs merged successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = "merged_document.pdf"
                        st.download_button(
                            label="📥 Download Merged PDF",
                            data=merged_pdf,
                            file_name=output_filename,
                            mime="application/pdf"
                        )
        elif uploaded_files and len(uploaded_files) == 1:
            st.warning("Please select at least two PDF files to merge.")
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "Split PDF":
        st.markdown('<h2 class="sub-header">Split PDF File</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file to split",
            type=["pdf"],
            help="Select a PDF file to split into multiple files"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Get page range
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            page_count = len(pdf_document)
            pdf_document.close()
            uploaded_file.seek(0)
            
            st.write(f"**Total pages:** {page_count}")
            
            split_option = st.radio(
                "Split by:",
                ["Single page", "Page range"]
            )
            
            if split_option == "Single page":
                # Convert button for single page split
                if st.button("Split PDF into Single Pages"):
                    with st.spinner("Splitting PDF into single pages..."):
                        # Split PDF
                        zip_data = split_pdf(uploaded_file, "single")
                        
                        if zip_data:
                            st.markdown('<div class="success-box">✅ PDF split successfully!</div>', unsafe_allow_html=True)
                            
                            # Download button
                            output_filename = "split_pages.zip"
                            st.download_button(
                                label="📥 Download Split Pages (ZIP)",
                                data=zip_data,
                                file_name=output_filename,
                                mime="application/zip"
                            )
            
            else:  # Page range
                col1, col2 = st.columns(2)
                with col1:
                    start_page = st.number_input("Start page", min_value=1, max_value=page_count, value=1)
                with col2:
                    end_page = st.number_input("End page", min_value=1, max_value=page_count, value=page_count)
                
                if start_page > end_page:
                    st.error("Start page cannot be greater than end page.")
                else:
                    # Convert button for page range split
                    if st.button("Split PDF by Page Range"):
                        with st.spinner("Splitting PDF by page range..."):
                            # Split PDF
                            pdf_data = split_pdf(uploaded_file, "range", start_page, end_page)
                            
                            if pdf_data:
                                st.markdown('<div class="success-box">✅ PDF split successfully!</div>', unsafe_allow_html=True)
                                
                                # Download button
                                output_filename = f"pages_{start_page}_to_{end_page}.pdf"
                                st.download_button(
                                    label="📥 Download PDF Extract",
                                    data=pdf_data,
                                    file_name=output_filename,
                                    mime="application/pdf"
                                )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "Compress PDF":
        st.markdown('<h2 class="sub-header">Compress PDF File</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file to compress",
            type=["pdf"],
            help="Select a PDF file to reduce its file size"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**Original file details:**")
            st.json(file_details)
            
            # Compression level
            compression_level = st.slider(
                "Compression level",
                min_value=1,
                max_value=5,
                value=3,
                help="Higher values mean more compression but potentially lower quality"
            )
            
            # Convert button
            if st.button("Compress PDF"):
                with st.spinner("Compressing PDF..."):
                    # Compress PDF
                    compressed_pdf = compress_pdf(uploaded_file, compression_level)
                    
                    if compressed_pdf:
                        original_size = uploaded_file.size / 1024
                        new_size = len(compressed_pdf.getvalue()) / 1024
                        reduction = ((original_size - new_size) / original_size) * 100
                        
                        st.markdown(f'<div class="success-box">✅ PDF compressed successfully! Size reduced from {original_size:.2f} KB to {new_size:.2f} KB ({reduction:.1f}% reduction)</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = "compressed_" + uploaded_file.name
                        st.download_button(
                            label="📥 Download Compressed PDF",
                            data=compressed_pdf,
                            file_name=output_filename,
                            mime="application/pdf"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "PDF to PowerPoint":
        st.markdown('<h2 class="sub-header">PDF to PowerPoint Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=["pdf"],
            help="Select a PDF file to convert to PowerPoint format"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Convert button
            if st.button("Convert PDF to PowerPoint"):
                with st.spinner("Converting PDF to PowerPoint..."):
                    # Convert PDF to PowerPoint
                    pptx_data = pdf_to_pptx(uploaded_file)
                    
                    if pptx_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = Path(uploaded_file.name).stem + ".pptx"
                        st.download_button(
                            label="📥 Download PowerPoint Presentation",
                            data=pptx_data,
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "PDF to JPG":
        st.markdown('<h2 class="sub-header">PDF to JPG Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=["pdf"],
            help="Select a PDF file to convert to JPG image"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Get page count
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            page_count = len(pdf_document)
            pdf_document.close()
            uploaded_file.seek(0)
            
            st.write(f"**Total pages:** {page_count}")
            
            if page_count > 1:
                page_number = st.number_input(
                    "Page to convert to JPG",
                    min_value=1,
                    max_value=page_count,
                    value=1,
                    help="Select which page of the PDF to convert to JPG"
                )
            else:
                page_number = 1
            
            # Convert button
            if st.button("Convert PDF to JPG"):
                with st.spinner("Converting PDF to JPG..."):
                    # Convert PDF to JPG
                    jpg_data = pdf_to_jpg(uploaded_file, page_number-1)
                    
                    if jpg_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Display image
                        st.image(jpg_data, caption=f"Page {page_number}", use_column_width=True)
                        
                        # Download button
                        output_filename = f"{Path(uploaded_file.name).stem}_page{page_number}.jpg"
                        st.download_button(
                            label="📥 Download JPG Image",
                            data=jpg_data,
                            file_name=output_filename,
                            mime="image/jpeg"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "JPG to PDF":
        st.markdown('<h2 class="sub-header">JPG to PDF Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_files = st.file_uploader(
            "Choose JPG images to convert to PDF",
            type=["jpg", "jpeg", "png"],
            help="Select one or more image files to convert to PDF",
            accept_multiple_files=True
        )
        
        if uploaded_files:
            st.write(f"**Files to convert:** {len(uploaded_files)} images selected")
            
            # Display image previews
            if len(uploaded_files) <= 5:  # Only show preview for small number of files
                cols = st.columns(min(3, len(uploaded_files)))
                for i, img_file in enumerate(uploaded_files):
                    with cols[i % 3]:
                        st.image(img_file, caption=img_file.name, use_column_width=True)
            
            # Convert button
            if st.button("Convert JPG to PDF"):
                with st.spinner("Converting images to PDF..."):
                    # Convert JPG to PDF
                    pdf_data = jpg_to_pdf(uploaded_files)
                    
                    if pdf_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = "converted_document.pdf"
                        st.download_button(
                            label="📥 Download PDF Document",
                            data=pdf_data,
                            file_name=output_filename,
                            mime="application/pdf"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "PDF to Excel":
        st.markdown('<h2 class="sub-header">PDF to Excel Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=["pdf"],
            help="Select a PDF file to convert to Excel format"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Convert button
            if st.button("Convert PDF to Excel"):
                with st.spinner("Converting PDF to Excel..."):
                    # Convert PDF to Excel
                    excel_data = pdf_to_excel(uploaded_file)
                    
                    if excel_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = Path(uploaded_file.name).stem + ".xlsx"
                        st.download_button(
                            label="📥 Download Excel Spreadsheet",
                            data=excel_data,
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        )
        st.markdown('</div>', unsafe_allow_html=True)
    
    elif tool_option == "Excel to PDF":
        st.markdown('<h2 class="sub-header">Excel to PDF Converter</h2>', unsafe_allow_html=True)
        st.markdown('<div class="converter-option">', unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader(
            "Choose an Excel file",
            type=["xlsx", "xls"],
            help="Select an Excel file to convert to PDF format"
        )
        
        if uploaded_file is not None:
            # Display file info
            file_details = {
                "Filename": uploaded_file.name,
                "File size": f"{uploaded_file.size / 1024:.2f} KB"
            }
            st.write("**File details:**")
            st.json(file_details)
            
            # Convert button
            if st.button("Convert Excel to PDF"):
                with st.spinner("Converting Excel to PDF..."):
                    # Convert Excel to PDF
                    pdf_data = excel_to_pdf(uploaded_file)
                    
                    if pdf_data:
                        st.markdown('<div class="success-box">✅ Conversion completed successfully!</div>', unsafe_allow_html=True)
                        
                        # Download button
                        output_filename = Path(uploaded_file.name).stem + ".pdf"
                        st.download_button(
                            label="📥 Download PDF Document",
                            data=pdf_data,
                            file_name=output_filename,
                            mime="application/pdf"
                        )
                    else:
                        st.info("💡 If conversion fails, please ensure LibreOffice is installed on your system.")
        st.markdown('</div>', unsafe_allow_html=True)
    
    # Footer
    st.markdown("---")
    st.markdown('<div class="footer">Made with Streamlit • Advanced Document Converter</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()